import os
import uuid
from typing import Optional, Dict, Any
from pathlib import Path
import asyncio
import aiofiles
import mimetypes

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import tiktoken
from PIL import Image

from config import settings

# Optional imports for OCR
try:
    if settings.ocr_enabled:
        import easyocr
        OCR_AVAILABLE = True
    else:
        OCR_AVAILABLE = False
except ImportError:
    OCR_AVAILABLE = False


class DocumentParser:
    def __init__(self):
        self.encoding = tiktoken.get_encoding("cl100k_base")  # GPT-4 encoding
        self.ocr_reader = None
        
        if settings.ocr_enabled and OCR_AVAILABLE:
            try:
                self.ocr_reader = easyocr.Reader(['en'])
            except Exception as e:
                print(f"Warning: OCR initialization failed: {e}")
    
    async def parse_document(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Parse document and extract text content"""
        try:
            file_ext = Path(filename).suffix.lower()
            
            if file_ext == '.pdf':
                content = await self._parse_pdf(file_path)
            elif file_ext == '.docx':
                content = await self._parse_docx(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                content = await self._parse_pptx(file_path)
            elif file_ext in ['.txt', '.md']:
                content = await self._parse_text(file_path)
            elif file_ext in ['.png', '.jpg', '.jpeg'] and self.ocr_reader:
                content = await self._parse_image(file_path)
            else:
                # Fallback to text parsing
                content = await self._parse_text_fallback(file_path)
            
            # Count tokens
            token_count = self._count_tokens(content)
            
            return {
                "raw_content": content,
                "parsed_content": content,
                "token_count": token_count,
                "file_type": file_ext,
                "metadata": {
                    "parsing_method": self._get_parsing_method(file_ext),
                    "original_filename": filename,
                    "ocr_enabled": settings.ocr_enabled
                }
            }
            
        except Exception as e:
            raise Exception(f"Failed to parse document: {str(e)}")
    
    async def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF using PyPDF2"""
        try:
            content = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            
            return content.strip()
        except Exception as e:
            raise Exception(f"PDF parsing failed: {str(e)}")
    
    async def _parse_docx(self, file_path: str) -> str:
        """Parse DOCX using python-docx"""
        try:
            doc = DocxDocument(file_path)
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            return content.strip()
        except Exception as e:
            raise Exception(f"DOCX parsing failed: {str(e)}")
    
    async def _parse_pptx(self, file_path: str) -> str:
        """Parse PPTX using python-pptx"""
        try:
            prs = Presentation(file_path)
            content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content.strip()
        except Exception as e:
            raise Exception(f"PPTX parsing failed: {str(e)}")
    
    async def _parse_text(self, file_path: str) -> str:
        """Parse text files"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                content = await file.read()
            return content
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                async with aiofiles.open(file_path, 'r', encoding='latin-1') as file:
                    content = await file.read()
                return content
            except Exception as e:
                raise Exception(f"Text parsing failed: {str(e)}")
    
    async def _parse_text_fallback(self, file_path: str) -> str:
        """Fallback text parsing for unsupported formats"""
        try:
            return await self._parse_text(file_path)
        except Exception:
            return f"Could not parse file: {os.path.basename(file_path)}"
    
    async def _parse_image(self, file_path: str) -> str:
        """Parse images using OCR"""
        if not self.ocr_reader:
            raise Exception("OCR is not available")
        
        try:
            # Run OCR in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                None, 
                self.ocr_reader.readtext, 
                file_path
            )
            
            # Extract text from OCR results
            content = ""
            for detection in result:
                content += detection[1] + " "
            
            return content.strip()
        except Exception as e:
            raise Exception(f"OCR failed: {str(e)}")
    
    def _count_tokens(self, text: str) -> int:
        """Count tokens in text using tiktoken"""
        try:
            tokens = self.encoding.encode(text)
            return len(tokens)
        except:
            # Fallback to word count * 1.3 (rough token approximation)
            return int(len(text.split()) * 1.3)
    
    def _get_parsing_method(self, file_ext: str) -> str:
        """Get the parsing method used for the file type"""
        method_map = {
            '.pdf': 'PyPDF2',
            '.docx': 'python-docx',
            '.pptx': 'python-pptx',
            '.ppt': 'python-pptx',
            '.txt': 'direct text',
            '.md': 'direct text',
            '.png': 'EasyOCR' if self.ocr_reader else 'not supported',
            '.jpg': 'EasyOCR' if self.ocr_reader else 'not supported',
            '.jpeg': 'EasyOCR' if self.ocr_reader else 'not supported'
        }
        return method_map.get(file_ext, 'text fallback')


class WebScraper:
    def __init__(self):
        pass
    
    async def scrape_url(self, url: str, include_links: bool = False) -> Dict[str, Any]:
        """Simple web scraping using httpx"""
        try:
            import httpx
            
            async with httpx.AsyncClient() as client:
                response = await client.get(url, timeout=30.0)
                response.raise_for_status()
                
                # Simple HTML text extraction
                content = response.text
                
                # Basic HTML tag removal (very simple)
                import re
                content = re.sub(r'<[^>]+>', ' ', content)
                content = re.sub(r'\s+', ' ', content).strip()
                
                # Count tokens
                encoding = tiktoken.get_encoding("cl100k_base")
                token_count = len(encoding.encode(content))
                
                return {
                    "raw_content": content,
                    "parsed_content": content,
                    "token_count": token_count,
                    "file_type": "web",
                    "metadata": {
                        "source_url": url,
                        "parsing_method": "Simple HTTP + regex",
                        "include_links": include_links
                    }
                }
                
        except Exception as e:
            raise Exception(f"Failed to scrape URL: {str(e)}")